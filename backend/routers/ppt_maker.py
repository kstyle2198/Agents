from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
import os
import tempfile
from typing import TypedDict, List
from pptx import Presentation
from pptx.util import Pt
from langgraph.graph import StateGraph, END

# LLM 활용을 위한 임포트 추가
from langchain_groq import ChatGroq
from langchain_core.messages import SystemMessage, HumanMessage

from utils.setlogger import setup_logger
logger = setup_logger(f"{__name__}")

from dotenv import load_dotenv
load_dotenv()
think_model = os.getenv("NO_THINK_MODEL")

# =========================================
# 1. 코어 로직: 서브 주제 및 핵심 문장 추출
# =========================================
def extract_subtopics_and_key_points(text: str) -> List[dict]:
    """
    LLM을 활용하여 텍스트에서 서브 주제와 주제별 핵심 문장 5~7개를 추출합니다.
    """
    try:
        llm = ChatGroq(model=think_model, temperature=0.5)
        
        messages = [
            SystemMessage(content=(
                "당신은 텍스트 분석 전문가입니다. "
                "주어진 텍스트를 분석하여 3~5개의 주요 서브 주제를 식별하고, "
                "각 서브 주제별로 5~7개의 핵심 문장을 추출하세요. "
                "응답은 다음 JSON 형식으로만 출력하세요:\n"
                "[\n"
                "  {\n"
                "    \"subtopic\": \"서브 주제 1\",\n"
                "    \"key_points\": [\"핵심 문장 1\", \"핵심 문장 2\", ...]\n"
                "  },\n"
                "  {\n"
                "    \"subtopic\": \"서브 주제 2\", \n"
                "    \"key_points\": [\"핵심 문장 1\", \"핵심 문장 2\", ...]\n"
                "  }\n"
                "]\n"
                "핵심 문장은 간결하고 명확하게 작성하며, 개조식 형식으로 제공하세요."
            )),
            HumanMessage(content=f"다음 텍스트를 분석해주세요:\n\n{text}")
        ]
        
        response = llm.invoke(messages)
        
        # JSON 응답 파싱
        import json
        try:
            subtopics_data = json.loads(response.content)
            logger.info(subtopics_data)
            return subtopics_data
        except json.JSONDecodeError:
            # JSON 파싱 실패 시 기본 구조 반환
            logger.error("서브 주제 및 문장 추출 에러")
            return [
                {
                    "subtopic": "주요 내용",
                    "key_points": ["텍스트 분석 중 오류가 발생했습니다.", "원본 텍스트를 확인해주세요."]
                }
            ]
            
    except Exception as e:
        raise RuntimeError(f"서브 주제 추출 중 오류 발생: {e}")

# =========================================
# 2. 코어 로직: PPTX 생성 함수 (Tool)
# =========================================
def create_pptx_file_with_subtopics(text: str) -> str:
    """
    주어진 텍스트를 LLM을 이용해 서브 주제별로 분석하고, 
    각 서브 주제별로 PowerPoint 슬라이드를 생성합니다.
    """
    try:
        # 서브 주제 및 핵심 문장 추출
        subtopics_data = extract_subtopics_and_key_points(text)
        
        # 프레젠테이션 객체 생성
        prs = Presentation()

        # 제목 슬라이드 추가
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = "텍스트 분석 결과"
        if hasattr(title_slide, 'placeholders') and len(title_slide.placeholders) > 1:
            subtitle = title_slide.placeholders[1]
            subtitle.text = f"총 {len(subtopics_data)}개의 서브 주제"

        # 각 서브 주제별로 슬라이드 추가
        for i, subtopic_data in enumerate(subtopics_data, 1):
            # 제목+내용 슬라이드 레이아웃 사용
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # 제목 및 내용 설정
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            # 서브 주제를 제목으로 설정
            title_shape.text = f"{i}. {subtopic_data['subtopic']}"
            
            # 핵심 문장들을 내용으로 추가
            tf = body_shape.text_frame
            tf.text = ""  # 초기 텍스트 초기화
            
            # 각 핵심 문장을 개별 paragraph로 추가
            for j, point in enumerate(subtopic_data['key_points']):
                if j == 0:
                    p = tf.paragraphs[0]
                    p.text = point
                else:
                    p = tf.add_paragraph()
                    p.text = point
                
                # 폰트 사이즈 설정
                for run in p.runs:
                    run.font.size = Pt(20)

        # 임시 파일로 저장
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            prs.save(tmp_file.name)
            return tmp_file.name
            
    except Exception as e:
        raise RuntimeError(f"PPTX 생성 프로세스 중 오류 발생: {e}")

# =========================================
# 3. LangGraph 워크플로우 정의
# =========================================

# 그래프의 상태(State) 정의
class GraphState(TypedDict):
    input_text: str  # 입력받은 텍스트
    pptx_path: str | None   # 생성된 PPTX 파일 경로 (None 허용)
    error: str | None       # 에러 메시지 (None 허용)
    subtopics: List[dict] | None  # 추출된 서브 주제 데이터 (None 허용)

# 노드 1: PPTX 생성 노드
def generate_pptx_node(state: GraphState) -> GraphState:
    try:
        text = state["input_text"]
        # 개선된 PPTX 생성 함수 호출 (서브 주제별 슬라이드 생성)
        file_path = create_pptx_file_with_subtopics(text)
        return {"pptx_path": file_path, "error": None}
    except Exception as e:
        return {"pptx_path": None, "error": str(e)}

# LangGraph 그래프 구성
workflow = StateGraph(GraphState)
workflow.add_node("generator", generate_pptx_node)

workflow.set_entry_point("generator")
workflow.add_edge("generator", END)
app_graph = workflow.compile()

# Router 구성
pptx_maker = APIRouter()

class PPTXRequest(BaseModel):
    text: str

class PPTXResponse(BaseModel):
    message: str
    subtopics_count: int
    file_path: str

@pptx_maker.post("/generate-pptx", tags=["PPTX_MAKER"], summary="텍스트 분석 및 다중 슬라이드 PPTX 생성")
async def generate_pptx_endpoint(request: PPTXRequest):
    """
    긴 텍스트를 입력받아 LLM으로 서브 주제를 추출하고, 
    각 서브 주제별로 PPTX 슬라이드를 생성하여 반환합니다.
    """
    inputs = {"input_text": request.text}
    result = app_graph.invoke(inputs)

    if result.get("error"):
        raise HTTPException(status_code=500, detail=result["error"])

    file_path = result.get("pptx_path")
    if file_path and os.path.exists(file_path):
        return FileResponse(
            path=file_path,
            filename="analyzed_presentation.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    else:
        raise HTTPException(status_code=500, detail="파일 생성에 실패했습니다.")

@pptx_maker.post("/analyze-text", tags=["PPTX_MAKER"], summary="텍스트 분석만 수행")
async def analyze_text_endpoint(request: PPTXRequest):
    """
    텍스트를 분석하여 서브 주제와 핵심 문장만 추출합니다.
    """
    try:
        subtopics_data = extract_subtopics_and_key_points(request.text)
        return {
            "message": "텍스트 분석이 완료되었습니다.",
            "subtopics_count": len(subtopics_data),
            "subtopics": subtopics_data
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"텍스트 분석 중 오류 발생: {e}")